from pathlib import Path
from pdf2image import convert_from_path
from pptx import Presentation
from io import BytesIO

pdf_file = "123.pdf"
print()
print("Converting file: " + pdf_file)
print()

# Prep presentation
prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]

# Create working folder
base_name = pdf_file.split(".pdf")[0]

# Convert PDF to list of images
print("Starting conversion...")
print(f"pdf_file: {pdf_file}")
print(f"pdf_file: {Path(pdf_file).absolute()}")
# slideimgs = convert_from_path(pdf_file, 300, fmt='ppm', thread_count=2)
slideimgs = convert_from_path(Path(pdf_file).absolute(), 300, fmt='ppm', thread_count=2)
print("...complete.")
print()

# Loop over slides
for i, slideimg in enumerate(slideimgs):
    if i % 10 == 0:
        print("Saving slide: " + str(i))

    imagefile = BytesIO()
    slideimg.save(imagefile, format='tiff')
    imagedata = imagefile.getvalue()
    imagefile.seek(0)
    width, height = slideimg.size

    # 调整幻灯片大小
    prs.slide_height = height * 9525
    prs.slide_width = width * 9525

    # Add slide
    slide = prs.slides.add_slide(blank_slide_layout)
    pic = slide.shapes.add_picture(imagefile, 0, 0, width=width * 9525, height=height * 9525)

# Save Powerpoint
print()
print("Saving file: " + base_name + ".pptx")
prs.save(base_name + '.pptx')
print("Conversion complete. :)")
print()
